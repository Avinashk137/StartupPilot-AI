import os
import io
import json
from typing import Optional
from fastapi import APIRouter, Depends, HTTPException, Query, Response
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from datetime import datetime, timezone
import structlog

from supabase import create_client, Client
from backend.core.dependencies import get_current_user
from backend.core.config import settings

logger = structlog.get_logger()
router = APIRouter(tags=["Exports"])

def get_supabase() -> Client:
    url = os.getenv("SUPABASE_URL", "")
    key = os.getenv("SUPABASE_SECRET_KEY", "")
    if not url or not key:
        raise HTTPException(status_code=500, detail="Missing Supabase credentials")
    return create_client(url, key)

@router.get("/stats")
async def get_export_stats(
    current_user = Depends(get_current_user),
    supabase: Client = Depends(get_supabase)
):
    try:
        user_id = current_user.id
        
        # We can't do complex aggregates easily via PostgREST without an RPC, 
        # so we'll fetch basic metrics by querying the table (only non-deleted).
        res = supabase.table("exports").select("status, download_count").eq("user_id", str(user_id)).eq("is_deleted", False).execute()
        
        records = res.data or []
        total = len(records)
        completed = sum(1 for r in records if r.get("status") == "completed")
        pending = sum(1 for r in records if r.get("status") in ["processing", "pending"])
        failed = sum(1 for r in records if r.get("status") == "failed")
        downloads = sum(r.get("download_count") or 0 for r in records)
        
        last_res = supabase.table("exports").select("created_at").eq("user_id", str(user_id)).eq("is_deleted", False).order("created_at", desc=True).limit(1).execute()
        last_date = last_res.data[0].get("created_at") if last_res.data else None
        
        return {
            "total_reports": total,
            "completed": completed,
            "pending": pending,
            "failed": failed,
            "total_downloads": downloads,
            "last_export_date": last_date
        }
    except Exception as e:
        logger.error(f"Failed to fetch export stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("")
@router.get("/")
async def list_exports(
    search: Optional[str] = None,
    status: Optional[str] = None,
    report_type: Optional[str] = None,
    sort: Optional[str] = "newest",
    page: int = 1,
    limit: int = 20,
    current_user = Depends(get_current_user),
    supabase: Client = Depends(get_supabase)
):
    try:
        user_id = current_user.id
        
        query = supabase.table("exports").select("id, project_id, project_name, industry, country, report_type, status, generated_by, version, download_count, created_at, updated_at", count="exact")
        query = query.eq("user_id", str(user_id)).eq("is_deleted", False)
        
        if status and status.lower() != "all":
            query = query.eq("status", status.lower())
        
        if report_type and report_type.lower() != "all":
            query = query.eq("report_type", report_type)
            
        if search:
            # Simple ilike on project_name
            query = query.ilike("project_name", f"%{search}%")
            
        if sort == "newest":
            query = query.order("created_at", desc=True)
        elif sort == "oldest":
            query = query.order("created_at", desc=False)
        elif sort == "downloads":
            query = query.order("download_count", desc=True)
        elif sort == "project":
            query = query.order("project_name", desc=False)
        elif sort == "industry":
            query = query.order("industry", desc=False)
        else:
            query = query.order("created_at", desc=True)
            
        start = (page - 1) * limit
        end = start + limit - 1
        query = query.range(start, end)
        
        res = query.execute()
        
        return {
            "data": res.data,
            "total": res.count if res.count else len(res.data),
            "page": page,
            "limit": limit
        }
    except Exception as e:
        logger.error(f"Failed to list exports: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{export_id}/preview")
async def preview_export(
    export_id: str,
    current_user = Depends(get_current_user),
    supabase: Client = Depends(get_supabase)
):
    try:
        res = supabase.table("exports").select("content, markdown, html").eq("id", export_id).eq("user_id", str(current_user.id)).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="Export not found")
        return res.data[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/{export_id}")
async def delete_export(
    export_id: str,
    current_user = Depends(get_current_user),
    supabase: Client = Depends(get_supabase)
):
    try:
        res = supabase.table("exports").update({"is_deleted": True}).eq("id", export_id).eq("user_id", str(current_user.id)).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="Export not found")
        return {"success": True}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/{export_id}/download/{format}")
async def download_export(
    export_id: str,
    format: str,
    current_user = Depends(get_current_user),
    supabase: Client = Depends(get_supabase)
):
    format = format.lower()
    if format not in ["docx", "csv", "pptx", "pdf"]:
        raise HTTPException(status_code=400, detail=f"Format {format} not supported")
        
    try:
        res = supabase.table("exports").select("*").eq("id", export_id).eq("user_id", str(current_user.id)).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="Export not found")
            
        export = res.data[0]
        content = export.get("content") or {}
        project_name = export.get("project_name", "Project")
        report_type = export.get("report_type", "report")
        
        filename = f"{project_name.replace(' ', '_')}_{report_type}.{format}"
        file_stream = io.BytesIO()
        
        if format == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            title = doc.add_heading(f"{project_name} - {report_type.replace('_', ' ').title()}", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d')}")
            
            def process_dict(d, level=1):
                for k, v in d.items():
                    if isinstance(v, dict):
                        doc.add_heading(k.replace('_', ' ').title(), level=level)
                        process_dict(v, level + 1)
                    elif isinstance(v, list):
                        doc.add_heading(k.replace('_', ' ').title(), level=level)
                        for item in v:
                            if isinstance(item, dict):
                                process_dict(item, level + 2)
                            else:
                                doc.add_paragraph(str(item), style='List Bullet')
                    else:
                        doc.add_heading(k.replace('_', ' ').title(), level=level)
                        doc.add_paragraph(str(v))
                        
            process_dict(content)
            doc.save(file_stream)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            
        elif format == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"{project_name} - {report_type.replace('_', ' ').title()}"
            subtitle.text = f"StartupPilot AI Export\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
            
            bullet_slide_layout = prs.slide_layouts[1]
            for key, value in content.items():
                if isinstance(value, str):
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    title_shape.text = key.replace('_', ' ').title()
                    tf = body_shape.text_frame
                    
                    sentences = [s.strip() for s in value.split('.') if len(s.strip()) > 5]
                    for i, s in enumerate(sentences[:6]):
                        if i == 0:
                            tf.text = s + "."
                        else:
                            p = tf.add_paragraph()
                            p.text = s + "."
            
            prs.save(file_stream)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            
        elif format == "csv":
            import csv
            file_stream = io.StringIO()
            writer = csv.writer(file_stream)
            writer.writerow(["Field", "Value"])
            
            def flatten_dict(d, parent_key='', sep='_'):
                items = []
                for k, v in d.items():
                    new_key = f"{parent_key}{sep}{k}" if parent_key else k
                    if isinstance(v, dict):
                        items.extend(flatten_dict(v, new_key, sep=sep).items())
                    elif isinstance(v, list):
                        items.append((new_key, ", ".join(str(i) for i in v)))
                    else:
                        items.append((new_key, str(v)))
                return dict(items)
                
            flat = flatten_dict(content)
            for k, v in flat.items():
                writer.writerow([k, v])
            
            file_stream.seek(0)
            csv_data = file_stream.getvalue()
            
            supabase.table("exports").update({
                "download_count": (export.get("download_count") or 0) + 1,
                "last_downloaded": datetime.now(timezone.utc).isoformat()
            }).eq("id", export_id).execute()
            
            return Response(content=csv_data, media_type="text/csv", headers={"Content-Disposition": f"attachment; filename={filename}"})
            
        elif format == "pdf":
            # Just return HTML or markdown, frontend will render PDF because pure python pdf requires heavy libs
            return Response(content=json.dumps({"success": True}), media_type="application/json")

        file_stream.seek(0)
        
        supabase.table("exports").update({
            "download_count": (export.get("download_count") or 0) + 1,
            "last_downloaded": datetime.now(timezone.utc).isoformat()
        }).eq("id", export_id).execute()
        
        return StreamingResponse(
            iter([file_stream.getvalue()]),
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}"
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Download failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))
